import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt

# Configuration
INPUT_FILE = "docs/notebook_llm/FINAL_THESIS_SOURCE_1_TO_80.md"
OUTPUT_FILE = "docs/notebook_llm/EcoSync_Final_Thesis_Presentation.pptx"

def clean_text(text):
    # Remove markdown bolding
    text = text.replace("**", "").replace("__", "")
    # Remove header markers if any remain
    text = re.sub(r'^#+\s*', '', text)
    return text.strip()

def extract_pages(filename):
    with open(filename, 'r', encoding='utf-8') as f:
        content = f.read()

    # IMPROVED REGEX:
    # 1. Matches any number of # at start of line
    # 2. Matches [PAGE X] or PAGE X (case insensitive)
    # 3. Captures Title until newline
    # 4. Captures Body using non-greedy match until next header or End of String
    
    # We use explicit capture groups for Page Num (1) and Title (2)
    # The Body is group 3.
    
    pattern = r'(?:^|\n)#{1,6}\s*(?:\[?PAGE\s*(\d+)\]?[:\s])(.*?)\n(.*?)(?=\n#{1,6}\s*\[?PAGE|\Z)'
    
    matches = re.findall(pattern, content, re.DOTALL | re.IGNORECASE)
    
    pages = []
    print(f"DEBUG: Found {len(matches)} raw matches via Regex.")
    
    for m in matches:
        page_num = m[0]
        title = m[1].replace("]", "").replace(":", "").strip()
        raw_body = m[2].strip()
        
        # Clean up body
        # Remove (Word Count...) lines
        body = re.sub(r'\*\*\(Word Count:.*?\)\*\*', '', raw_body)
        body = body.replace("---", "").strip()
        
        pages.append({
            "page": int(page_num),
            "title": clean_text(title),
            "content": body
        })
        
    return pages

def create_presentation(pages):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "EcoSync: An AI-Driven Sustainable Cloud Data Center Simulator"
    subtitle.text = "Final Year Project Thesis Presentation\nGenerated from 80-Page Source"

    print(f"🚀 Processing {len(pages)} pages into slides...")

    for page_data in pages: # RENAMED variable to avoid shadowing
        # Create standard slide
        bullet_slide_layout = prs.slide_layouts[1] # Title and Body
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # Set Title
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = f"{page_data['page']}. {page_data['title']}"
        
        # Set Body (Bullet Points)
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # Clear default prompt text
        
        # Heuristic Summarizer
        paragraphs = [para.strip() for para in page_data['content'].split('\n') if para.strip()]
        
        bullet_count = 0
        for para in paragraphs:
            # Skip code blocks and meaningless lines
            if para.startswith("```") or para.startswith("import ") or para.startswith("def ") or len(para) < 20:
                continue
            
            clean_para = clean_text(para)
            
            # Simple Heuristic: Take the first sentence of paragraphs
            first_sentence = clean_para.split('. ')[0] + "."
            
            if len(first_sentence) > 10:
                p_obj = tf.add_paragraph() # RENAMED variable
                p_obj.text = first_sentence[:120] + "..." if len(first_sentence) > 120 else first_sentence
                p_obj.level = 0
                bullet_count += 1
            
            if bullet_count >= 4: # Max 4 bullets
                break
                
        # Set Speaker Notes (FULL CONTENT)
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = page_data['content'] # Full text goes here
        
    prs.save(OUTPUT_FILE)
    print(f"✅ Successfully created {OUTPUT_FILE} with {len(prs.slides)} slides.")

if __name__ == "__main__":
    if not os.path.exists(INPUT_FILE):
        print(f"❌ Input file not found: {INPUT_FILE}")
    else:
        pages = extract_pages(INPUT_FILE)
        if pages:
            # Sort pages just in case regex found them out of order (unlikely but safe)
            pages.sort(key=lambda x: x['page'])
            create_presentation(pages)
        else:
            print("❌ No pages found. Check regex patterns.")
